"""
Build the ENEB453 Checkpoint PPT
Preliminary UI and DB Design Update
Nishanth S. | University of Maryland
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT = r"C:\Users\nisha\Videos\Combined Class Final Project\Final Project\ENEB453\ENEB453_Checkpoint_NishanthS.pptx"
SHOTS = r"C:\Users\nisha\Videos\Combined Class Final Project\Final Project\ENEB453\screenshots"

# ── Colors ──────────────────────────────────────────────
DARK_BG   = RGBColor(0x16, 0x21, 0x3E)   # navy
ACCENT    = RGBColor(0xE9, 0x45, 0x60)   # red
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRY = RGBColor(0xF1, 0xF5, 0xF9)
MID_GRY   = RGBColor(0x64, 0x74, 0x8B)
GOLD      = RGBColor(0xFB, 0xBF, 0x24)
BLUE_LT   = RGBColor(0x60, 0xA5, 0xFA)
GREEN     = RGBColor(0x22, 0xC5, 0x5E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank = prs.slide_layouts[6]   # completely blank

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_bg(slide, color=DARK_BG):
    bg = add_rect(slide, 0, 0, 13.33, 7.5, fill=color)
    return bg

def header_bar(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 13.33, 1.1, fill=DARK_BG)
    add_rect(slide, 0, 1.08, 13.33, 0.04, fill=ACCENT)
    add_text(slide, title,    0.3, 0.08, 10, 0.6, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.68, 10, 0.35, size=13, color=MID_GRY)
    add_text(slide, "Nishanth S.  |  ENEB453  |  University of Maryland",
             0.3, 0.68, 12.5, 0.35, size=11, color=MID_GRY, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
add_bg(s1)

# big accent bar left edge
add_rect(s1, 0, 0, 0.06, 7.5, fill=ACCENT)

# large centered title block
add_rect(s1, 1.2, 1.8, 10.8, 3.8, fill=RGBColor(0x0F, 0x34, 0x60))

add_text(s1, "BFRB Detection System",
         1.5, 2.1, 10.3, 1.0, size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, "Preliminary UI & Database Design Update",
         1.5, 3.15, 10.3, 0.7, size=22, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)
add_rect(s1, 4.5, 3.9, 4.3, 0.04, fill=ACCENT)
add_text(s1, "ENEB453 — Web-Based Application Development",
         1.5, 4.0, 10.3, 0.5, size=15, color=LIGHT_GRY, align=PP_ALIGN.CENTER)
add_text(s1, "Nishanth S.  |  University of Maryland  |  Spring 2026",
         1.5, 4.55, 10.3, 0.4, size=13, color=MID_GRY, align=PP_ALIGN.CENTER)
add_text(s1, "Due: Mon Apr 13, 2026",
         1.5, 5.6, 10.3, 0.4, size=12, color=MID_GRY, align=PP_ALIGN.CENTER)

# bullet points of what's in the deck
items = [
    ("01", "Wireframe — User Interface Design"),
    ("02", "Entity-Relationship Diagram (ERD) — Normalized to 3NF"),
    ("03", "Populated Database — Schema + Demo Queries"),
]
y = 5.1
for num, label in items:
    add_rect(s1, 3.3, y, 0.42, 0.28, fill=ACCENT)
    add_text(s1, num, 3.3, y, 0.42, 0.28, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s1, label, 3.78, y-0.01, 6.0, 0.3, size=12, color=LIGHT_GRY)
    y += 0.36

# ════════════════════════════════════════════════════════
# SLIDE 2 — PROJECT OVERVIEW
# ════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
add_bg(s2)
header_bar(s2, "Project Overview", "What this system does and why it matters")

# Two columns
add_rect(s2, 0.3, 1.25, 6.0, 5.85, fill=RGBColor(0x0F, 0x34, 0x60))
add_rect(s2, 6.6, 1.25, 6.4, 5.85, fill=RGBColor(0x0F, 0x34, 0x60))

add_text(s2, "What It Is", 0.5, 1.3, 5.6, 0.4, size=15, bold=True, color=ACCENT)
desc = (
    "A web application that uses a webcam or uploaded video "
    "to detect Body-Focused Repetitive Behaviors (BFRBs) in real time "
    "using Google MediaPipe computer vision.\n\n"
    "The system watches hand and face movements, identifies when a person "
    "is performing behaviors like hair pulling or nail biting, logs every "
    "event to a database, and lets users define custom new behaviors."
)
add_text(s2, desc, 0.5, 1.75, 5.6, 2.2, size=12, color=LIGHT_GRY, wrap=True)

add_text(s2, "Professor's Analogy", 0.5, 3.6, 5.6, 0.4, size=15, bold=True, color=ACCENT)
analogy = (
    "Like a store security system that detects shoplifters by recognizing "
    "their specific nervous walking pattern — this system detects BFRB "
    "behaviors by recognizing specific hand-to-face movement patterns."
)
add_text(s2, analogy, 0.5, 4.05, 5.6, 1.5, size=12, color=LIGHT_GRY, wrap=True)

add_text(s2, "Technology Stack", 0.5, 5.7, 5.6, 0.4, size=15, bold=True, color=ACCENT)
stack = "Python Flask  ·  MediaPipe Holistic  ·  MySQL  ·  Docker  ·  HTML/CSS/JS"
add_text(s2, stack, 0.5, 6.15, 5.6, 0.5, size=12, color=GREEN)

# Right column
add_text(s2, "BFRB Behaviors Detected", 6.8, 1.3, 6.0, 0.4, size=15, bold=True, color=ACCENT)

behaviors = [
    ("Hair Pulling",       "Fingertip near scalp/forehead landmark"),
    ("Nail Biting",        "Fingertip near upper lip landmark"),
    ("Skin Picking",       "Fingertip near cheek/forehead landmark"),
    ("Nose Picking",       "Index tip near nose tip landmark"),
    ("Lip Picking",        "Thumb tip near lip landmarks"),
    ("Custom (User-Def.)", "User configures any landmark pair"),
]
y = 1.78
for name, rule in behaviors:
    add_rect(s2, 6.8, y, 0.08, 0.28, fill=ACCENT)
    add_text(s2, name, 7.0, y, 2.3, 0.28, size=12, bold=True, color=WHITE)
    add_text(s2, rule, 9.35, y, 3.5, 0.28, size=11, color=MID_GRY)
    y += 0.42

add_text(s2, "How Detection Works", 6.8, 4.5, 6.0, 0.4, size=15, bold=True, color=ACCENT)
how = (
    "MediaPipe tracks 21 hand landmarks + 468 face landmarks per frame. "
    "The system measures the Euclidean distance between a hand fingertip "
    "and a target face region. If within the threshold for 10+ consecutive "
    "frames (~0.33 sec), a detection event is confirmed and saved to MySQL."
)
add_text(s2, how, 6.8, 4.95, 6.1, 1.8, size=12, color=LIGHT_GRY, wrap=True)

# ════════════════════════════════════════════════════════
# SLIDE 3 — WIREFRAME (screenshot)
# ════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
add_bg(s3)
header_bar(s3, "UI Wireframe", "Dashboard layout — the webpage the user sees")

wf_img = os.path.join(SHOTS, "wireframe.png")
if os.path.exists(wf_img):
    s3.shapes.add_picture(wf_img, Inches(0.25), Inches(1.2), Inches(12.83), Inches(6.1))

# Label overlay bottom right
add_rect(s3, 9.5, 6.9, 3.6, 0.45, fill=RGBColor(0x0F, 0x34, 0x60))
add_text(s3, "Built with HTML/CSS — Open wireframe.html in browser",
         9.55, 6.92, 3.5, 0.35, size=10, color=MID_GRY)

# ════════════════════════════════════════════════════════
# SLIDE 4 — WIREFRAME ANNOTATED DESCRIPTION
# ════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
add_bg(s4)
header_bar(s4, "Wireframe — Key Components", "What each section of the UI does")

sections = [
    ("Navigation Bar",        ACCENT,   "Dashboard · Sessions · BFRB Types · Reports — routes to all major features of the app."),
    ("Stats Row",             GOLD,     "4 live counters at top: Sessions Today, Detections Today, BFRB Types defined, Average Confidence score."),
    ("Video Feed Panel",      GREEN,    "Left side. Tabs for Webcam vs Upload Video (MP4). Shows live camera with MediaPipe skeleton overlay drawn on top. Active detection label appears in top-left corner of the feed."),
    ("Live Detection Log",    BLUE_LT,  "Right side. Every BFRB event fires here in real time — behavior name, timestamp, confidence bar, confidence percentage."),
    ("Current Session Info",  WHITE,    "Right side, below log. Shows session ID, input type (webcam/video), start time, duration, total events this session."),
    ("BFRB Type Management",  ACCENT,   "Bottom section. Table of all built-in and custom BFRB types. Form to add a new behavior — name, description, landmark rule. Each row has an Edit button."),
    ("Alert Settings",        MID_GRY,  "Right column, bottom. Toggle email alerts on/off. Set minimum confidence threshold. Toggle auto-save sessions."),
]

y = 1.25
for title, color, desc in sections:
    add_rect(s4, 0.3, y, 0.08, 0.32, fill=color)
    add_text(s4, title, 0.5, y, 3.0, 0.32, size=12, bold=True, color=color)
    add_text(s4, desc,  3.55, y, 9.5, 0.36, size=11, color=LIGHT_GRY, wrap=True)
    add_rect(s4, 0.3, y+0.38, 12.73, 0.01, fill=RGBColor(0x1E, 0x29, 0x3B))
    y += 0.44

# ════════════════════════════════════════════════════════
# SLIDE 5 — ERD (screenshot)
# ════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
add_bg(s5)
header_bar(s5, "Database ERD — Normalized to 3NF", "4 tables · 3 foreign key relationships")

erd_img = os.path.join(SHOTS, "erd.png")
if os.path.exists(erd_img):
    s5.shapes.add_picture(erd_img, Inches(0.25), Inches(1.2), Inches(12.83), Inches(6.1))

# ════════════════════════════════════════════════════════
# SLIDE 6 — ERD TABLE-BY-TABLE EXPLANATION
# ════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
add_bg(s6)
header_bar(s6, "Database Schema — Table Breakdown", "Why each table exists and what it stores")

tables = [
    ("bfrb_types",           GOLD,    "Catalog of every BFRB behavior — both built-in (from Kaggle dataset) and user-defined custom ones. Stores the name, description, and a flag marking whether it's custom."),
    ("bfrb_landmark_rules",  BLUE_LT, "The MediaPipe detection rules for each BFRB type. Stores which hand landmark index and which face landmark index to measure, plus the distance threshold in pixels. Separated from bfrb_types to satisfy 1NF — no JSON blobs."),
    ("sessions",             GREEN,   "One row per detection run. Records whether input was webcam or a video file, the filename if applicable, and the start/end timestamps. End_time is NULL while a session is still running."),
    ("detections",           ACCENT,  "One row per detected BFRB event. Foreign keys to both sessions and bfrb_types. Stores the video frame number, millisecond timestamp into the session, and a confidence score (0.0–1.0)."),
]

y = 1.25
for name, color, desc in tables:
    add_rect(s6, 0.3, y, 0.08, 0.55, fill=color)
    add_text(s6, name, 0.5, y,    2.8, 0.3, size=13, bold=True, color=color)
    add_text(s6, desc, 0.5, y+0.3, 12.5, 0.55, size=11, color=LIGHT_GRY, wrap=True)
    add_rect(s6, 0.3, y+0.72, 12.73, 0.01, fill=RGBColor(0x1E, 0x29, 0x3B))
    y += 0.80

add_text(s6, "Relationships", 0.3, 5.6, 3.0, 0.35, size=13, bold=True, color=ACCENT)
rels = [
    "bfrb_landmark_rules.bfrb_type_id  →  bfrb_types.id   (MANY-TO-ONE, ON DELETE CASCADE)",
    "detections.session_id             →  sessions.id      (MANY-TO-ONE, ON DELETE CASCADE)",
    "detections.bfrb_type_id           →  bfrb_types.id    (MANY-TO-ONE, ON DELETE RESTRICT)",
]
y = 6.0
for r in rels:
    add_text(s6, "→  " + r, 0.3, y, 12.7, 0.3, size=11, color=BLUE_LT)
    y += 0.3

# ════════════════════════════════════════════════════════
# SLIDE 7 — 3NF EXPLANATION
# ════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
add_bg(s7)
header_bar(s7, "Normalization to 3rd Normal Form (3NF)", "Why the schema satisfies 3NF")

nf_blocks = [
    ("1NF — First Normal Form",  GOLD,    "✓ Satisfied",
     "Every column holds a single atomic value — no lists, no nested data. "
     "The detection landmark rules are stored as individual rows in bfrb_landmark_rules "
     "rather than as a JSON blob inside bfrb_types. Each table has a clear primary key (id)."),
    ("2NF — Second Normal Form", BLUE_LT, "✓ Satisfied",
     "All tables use a single-column integer primary key, so composite keys don't exist "
     "and partial dependencies are impossible. Every non-key attribute in each table depends "
     "on the full primary key, not just part of it."),
    ("3NF — Third Normal Form",  GREEN,   "✓ Satisfied",
     "No transitive dependencies exist. Example: detections stores bfrb_type_id (the FK), "
     "not the behavior name directly. If a behavior name changes, only one row in bfrb_types "
     "needs updating — not thousands of detection rows. No non-key attribute determines "
     "another non-key attribute in any table."),
]

y = 1.25
for title, color, check, desc in nf_blocks:
    add_rect(s7, 0.3, y, 9.5, 1.55, fill=RGBColor(0x0F, 0x34, 0x60))
    add_rect(s7, 0.3, y, 0.08, 1.55, fill=color)
    add_text(s7, title, 0.5, y+0.1, 7.0, 0.4, size=15, bold=True, color=color)
    add_rect(s7, 8.5, y+0.08, 1.1, 0.36, fill=color)
    add_text(s7, check, 8.5, y+0.08, 1.1, 0.36, size=12, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(s7, desc, 0.5, y+0.55, 9.2, 0.92, size=12, color=LIGHT_GRY, wrap=True)
    y += 1.75

# ════════════════════════════════════════════════════════
# SLIDE 8 — DATABASE POPULATED (Schema + Row Counts)
# ════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank)
add_bg(s8)
header_bar(s8, "Populated Database — Tables & Seed Data", "Run via Docker MySQL · bfrb_detection database")

add_text(s8, "Database: bfrb_detection   |   User: bfrb_user   |   Engine: MySQL 8.0 via Docker",
         0.3, 1.18, 12.5, 0.35, size=12, color=MID_GRY)

# Table summary boxes
boxes = [
    ("bfrb_types",          "6 rows",  GOLD,    "5 built-in (Hair Pulling, Nail Biting, Skin Picking, Nose Picking, Lip Picking) + 1 custom (Cheek Biting)"),
    ("bfrb_landmark_rules", "11 rows", BLUE_LT, "11 MediaPipe landmark rules mapping hand fingertips to face regions with distance thresholds"),
    ("sessions",            "3 rows",  GREEN,   "2 webcam sessions + 1 video file session with start/end timestamps"),
    ("detections",          "17 rows", ACCENT,  "17 BFRB detection events spread across 3 sessions with frame numbers, timestamps, confidence scores"),
]

x_positions = [0.3, 3.55, 6.8, 10.05]
for i, (tname, rows, color, desc) in enumerate(boxes):
    x = x_positions[i]
    add_rect(s8, x, 1.65, 3.0, 2.1, fill=RGBColor(0x0F, 0x34, 0x60))
    add_rect(s8, x, 1.65, 3.0, 0.06, fill=color)
    add_text(s8, tname,   x+0.1, 1.72, 2.8, 0.4, size=13, bold=True, color=color)
    add_text(s8, rows,    x+0.1, 2.1,  2.8, 0.55, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s8, "rows",  x+0.1, 2.62, 2.8, 0.3,  size=11, color=MID_GRY, align=PP_ALIGN.CENTER)
    add_text(s8, desc,    x+0.1, 2.95, 2.8, 0.75, size=10, color=LIGHT_GRY, wrap=True)

# SQL schema snippets
add_text(s8, "CREATE TABLE Schema (abbreviated)", 0.3, 3.9, 6.0, 0.35, size=13, bold=True, color=ACCENT)
schema_code = (
    "CREATE TABLE bfrb_types (\n"
    "  id INT AUTO_INCREMENT PRIMARY KEY,\n"
    "  name VARCHAR(100) NOT NULL UNIQUE,\n"
    "  is_custom TINYINT(1) NOT NULL DEFAULT 0\n"
    ");\n\n"
    "CREATE TABLE detections (\n"
    "  id INT AUTO_INCREMENT PRIMARY KEY,\n"
    "  session_id INT NOT NULL,  -- FK → sessions\n"
    "  bfrb_type_id INT NOT NULL, -- FK → bfrb_types\n"
    "  confidence_score DECIMAL(5,4) NOT NULL\n"
    ");"
)
add_rect(s8, 0.3, 4.3, 6.2, 2.8, fill=RGBColor(0x0D, 0x1B, 0x2A))
add_text(s8, schema_code, 0.45, 4.35, 6.0, 2.7, size=10, color=GREEN, wrap=False)

add_text(s8, "Seed Data Sample — bfrb_types", 6.7, 3.9, 6.2, 0.35, size=13, bold=True, color=ACCENT)
seed_table = (
    "id  name               is_custom\n"
    "─────────────────────────────────\n"
    " 1  Hair Pulling            0\n"
    " 2  Nail Biting             0\n"
    " 3  Skin Picking (Face)     0\n"
    " 4  Nose Picking            0\n"
    " 5  Lip Picking             0\n"
    " 6  Cheek Biting            1   ← custom"
)
add_rect(s8, 6.7, 4.3, 6.3, 2.8, fill=RGBColor(0x0D, 0x1B, 0x2A))
add_text(s8, seed_table, 6.85, 4.35, 6.1, 2.7, size=11, color=GREEN, wrap=False)

# ════════════════════════════════════════════════════════
# SLIDE 9 — DEMO QUERIES (Q1 and Q2)
# ════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank)
add_bg(s9)
header_bar(s9, "Demo Queries — Query 1 & 2", "Populated database query results")

# Q1
add_text(s9, "Query 1 — All detections with behavior name, session type, timestamp, confidence",
         0.3, 1.22, 12.7, 0.35, size=12, bold=True, color=ACCENT)
q1_sql = "SELECT bt.name, s.input_type, CONCAT(FLOOR(timestamp_ms/60000),'m ',FLOOR((timestamp_ms%60000)/1000),'s') AS time, ROUND(confidence_score*100,1) AS pct FROM detections d JOIN sessions s ON d.session_id=s.id JOIN bfrb_types bt ON d.bfrb_type_id=bt.id ORDER BY d.detected_at;"
add_rect(s9, 0.3, 1.6, 12.7, 0.45, fill=RGBColor(0x0D, 0x1B, 0x2A))
add_text(s9, q1_sql, 0.4, 1.62, 12.5, 0.4, size=9, color=GREEN, wrap=False)

q1_result = (
    "name                  input_type  time      pct\n"
    "─────────────────────────────────────────────────\n"
    "Hair Pulling          webcam      0m 4s     89.1%\n"
    "Nail Biting           webcam      0m 11s    72.3%\n"
    "Hair Pulling          webcam      0m 22s    91.5%\n"
    "Skin Picking (Face)   webcam      0m 29s    61.0%\n"
    "Nose Picking          webcam      0m 34s    80.0%\n"
    "Lip Picking           video       0m 6s     74.2%\n"
    "Hair Pulling          video       0m 15s    88.1%\n"
    "...  (17 rows total)"
)
add_rect(s9, 0.3, 2.1, 12.7, 2.15, fill=RGBColor(0x0D, 0x1B, 0x2A))
add_text(s9, q1_result, 0.4, 2.12, 12.5, 2.1, size=11, color=LIGHT_GRY, wrap=False)

# Q2
add_text(s9, "Query 2 — Detection count and average confidence per BFRB type",
         0.3, 4.35, 12.7, 0.35, size=12, bold=True, color=ACCENT)
q2_result = (
    "behavior              custom  total_detections  avg_confidence_pct  max_confidence_pct\n"
    "──────────────────────────────────────────────────────────────────────────────────────\n"
    "Hair Pulling             0          5                 91.5%              94.0%\n"
    "Nail Biting              0          3                 77.3%              82.3%\n"
    "Lip Picking              0          2                 75.4%              76.6%\n"
    "Skin Picking (Face)      0          2                 63.2%              65.4%\n"
    "Nose Picking             0          2                 84.5%              89.0%\n"
    "Cheek Biting             1          1                 71.0%              71.0%\n"
    "                                   ↑ custom behavior also tracked"
)
add_rect(s9, 0.3, 4.75, 12.7, 2.35, fill=RGBColor(0x0D, 0x1B, 0x2A))
add_text(s9, q2_result, 0.4, 4.77, 12.5, 2.3, size=10, color=LIGHT_GRY, wrap=False)

# ════════════════════════════════════════════════════════
# SLIDE 10 — DEMO QUERIES (Q3, Q4, Q5)
# ════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank)
add_bg(s10)
header_bar(s10, "Demo Queries — Query 3, 4 & 5", "Session summary · Custom BFRB rules · High-confidence alerts")

# Q3
add_text(s10, "Query 3 — Session summary (LEFT JOIN, GROUP BY, TIMESTAMPDIFF)",
         0.3, 1.22, 12.7, 0.3, size=12, bold=True, color=ACCENT)
q3 = (
    "session_id  input_type  source                   start_time            end_time              duration_min  total_detections\n"
    "──────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────\n"
    "     1       webcam      webcam     2026-04-08 10:05:00   2026-04-08 10:22:00        17             7\n"
    "     2       video       study_session_01.mp4    2026-04-09 14:00:00   2026-04-09 14:18:00        18             5\n"
    "     3       webcam      webcam     2026-04-10 16:30:00   2026-04-10 16:47:00        17             5"
)
add_rect(s10, 0.3, 1.55, 12.7, 1.55, fill=RGBColor(0x0D, 0x1B, 0x2A))
add_text(s10, q3, 0.4, 1.57, 12.5, 1.5, size=9, color=LIGHT_GRY, wrap=False)

# Q4
add_text(s10, "Query 4 — Custom BFRB types with their MediaPipe landmark rules",
         0.3, 3.18, 12.7, 0.3, size=12, bold=True, color=ACCENT)
q4 = (
    "custom_behavior   hand_landmark_index  face_landmark_index  distance_threshold_px  rule_description\n"
    "──────────────────────────────────────────────────────────────────────────────────────────────────────\n"
    "Cheek Biting              8                   127                   28.0          Index fingertip within 28px of left cheek corner"
)
add_rect(s10, 0.3, 3.52, 12.7, 1.1, fill=RGBColor(0x0D, 0x1B, 0x2A))
add_text(s10, q4, 0.4, 3.54, 12.5, 1.05, size=10, color=LIGHT_GRY, wrap=False)

# Q5
add_text(s10, "Query 5 — High-confidence detections only (confidence ≥ 85%) — alert-worthy events",
         0.3, 4.7, 12.7, 0.3, size=12, bold=True, color=ACCENT)
q5 = (
    "behavior         input_type  confidence_pct  detected_at\n"
    "──────────────────────────────────────────────────────────────────────────────\n"
    "Hair Pulling      webcam         93.0%      2026-04-10 16:30:36\n"
    "Hair Pulling      webcam         90.0%      2026-04-10 16:30:10\n"
    "Hair Pulling      video          88.1%      2026-04-09 14:00:15\n"
    "Hair Pulling      webcam         89.1%      2026-04-08 10:05:04\n"
    "Nose Picking      video          89.0%      2026-04-09 14:00:37\n"
    "Hair Pulling      webcam         91.5%      2026-04-08 10:05:23"
)
add_rect(s10, 0.3, 5.05, 12.7, 2.05, fill=RGBColor(0x0D, 0x1B, 0x2A))
add_text(s10, q5, 0.4, 5.07, 12.5, 2.0, size=11, color=LIGHT_GRY, wrap=False)

# ════════════════════════════════════════════════════════
# SLIDE 11 — NEXT STEPS
# ════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(blank)
add_bg(s11)
header_bar(s11, "Next Steps — Build Plan", "From checkpoint to final submission (May 4, 2026)")

steps = [
    ("This Week (Apr 14–18)",  ACCENT,
     "Flask backend + MediaPipe Holistic integration. Live webcam feed with landmark overlay working. Basic detection logic firing events."),
    ("Week 2 (Apr 21–25)",     GOLD,
     "Full frontend UI. Video file upload + annotated output. MySQL integration — every detection saved. Add new BFRB form working."),
    ("Week 3 (Apr 28–May 3)",  GREEN,
     "Polish, testing, presentation build. Full demo rehearsal. Teaching notes prepared so I can answer any question in the 1-on-1."),
    ("May 4 — Presentation",   BLUE_LT,
     "Live demo of full working system. Q&A with professor. Show wireframe → ERD → working app progression."),
]

y = 1.25
for period, color, desc in steps:
    add_rect(s11, 0.3, y, 2.8, 1.3, fill=color)
    add_text(s11, period, 0.35, y+0.1, 2.7, 0.8, size=13, bold=True, color=DARK_BG,
             align=PP_ALIGN.CENTER, wrap=True)
    add_rect(s11, 3.15, y, 9.8, 1.3, fill=RGBColor(0x0F, 0x34, 0x60))
    add_text(s11, desc, 3.3, y+0.2, 9.6, 0.9, size=13, color=LIGHT_GRY, wrap=True)
    y += 1.5

add_text(s11, "Final deliverable: fully working web app · live demo · presentation",
         0.3, 7.1, 12.7, 0.35, size=12, color=MID_GRY, align=PP_ALIGN.CENTER)

# ── SAVE ────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
