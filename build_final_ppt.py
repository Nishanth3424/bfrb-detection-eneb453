"""
Build the ENEB453 Final Project Defense Presentation (25 minutes)
BFRB Detection System - Nishanth S. | University of Maryland
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "ENEB453_Final_Defense_NishanthS.pptx")

# ── Colors (White + Dark Blue Professional Theme) ──
NAVY      = RGBColor(0x1A, 0x3A, 0x5C)
NAVY_DARK = RGBColor(0x0F, 0x25, 0x3D)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)
TEXT_DARK = RGBColor(0x1A, 0x23, 0x32)
TEXT_DIM  = RGBColor(0x5A, 0x6A, 0x7E)
ACCENT_RED = RGBColor(0xD9, 0x53, 0x4F)
GREEN     = RGBColor(0x2E, 0x8B, 0x57)
GOLD      = RGBColor(0xD4, 0x88, 0x0F)
BLUE_LT   = RGBColor(0x2A, 0x6C, 0xB6)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H


# ── Helpers ──────────────────────────────────────────────
def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, w, h, text, size=18, color=WHITE,
                 bold=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_slide(title, bullets, speaker_notes="", bg_color=NAVY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, bg_color)
    # Top accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_RED)
    # Title
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 title, size=36, color=WHITE, bold=True)
    # Underline
    add_rect(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), ACCENT_RED)
    # Bullets
    y = Inches(1.6)
    for bullet in bullets:
        if bullet.startswith("**"):
            # Sub-header
            clean = bullet.strip("*")
            add_text_box(slide, Inches(1.0), y, Inches(10.5), Inches(0.5),
                         clean, size=20, color=GOLD, bold=True)
            y += Inches(0.45)
        else:
            add_text_box(slide, Inches(1.2), y, Inches(10.5), Inches(0.45),
                         bullet, size=18, color=WHITE)
            y += Inches(0.42)
    # Slide number
    add_text_box(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.3),
                 str(len(prs.slides)), size=10, color=TEXT_DIM,
                 align=PP_ALIGN.RIGHT)
    # Speaker notes
    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes
    return slide


def add_two_col_slide(title, left_items, right_items, speaker_notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_RED)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 title, size=36, color=WHITE, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04), ACCENT_RED)
    # Left column
    y = Inches(1.6)
    for item in left_items:
        add_text_box(slide, Inches(0.8), y, Inches(5.5), Inches(0.42),
                     item, size=17, color=WHITE)
        y += Inches(0.40)
    # Right column
    y = Inches(1.6)
    for item in right_items:
        add_text_box(slide, Inches(6.8), y, Inches(5.5), Inches(0.42),
                     item, size=17, color=WHITE)
        y += Inches(0.40)
    add_text_box(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.3),
                 str(len(prs.slides)), size=10, color=TEXT_DIM,
                 align=PP_ALIGN.RIGHT)
    if speaker_notes:
        slide.notes_slide.notes_text_frame.text = speaker_notes
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY_DARK)
add_rect(slide, Inches(0), Inches(3.2), SLIDE_W, Inches(0.06), ACCENT_RED)
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "BFRB Detection System", size=52, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.6),
             "Real-Time Webcam-Based Body-Focused Repetitive Behavior Detection",
             size=22, color=LIGHT_BG, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
             "ENEB453 Final Project Defense", size=20, color=GOLD,
             bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.5),
             "Nishanth S.  |  University of Maryland  |  Spring 2026",
             size=18, color=TEXT_DIM, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.3), Inches(11), Inches(0.5),
             "May 11, 2026", size=16, color=TEXT_DIM, align=PP_ALIGN.CENTER)
slide.notes_slide.notes_text_frame.text = (
    "Good afternoon everyone, and thank you for being here. My name is Nishanth, "
    "and today I will be presenting my ENEB453 final project: a real-time "
    "Body-Focused Repetitive Behavior detection system.\n\n"
    "BFRBs — which stands for Body-Focused Repetitive Behaviors — are compulsive "
    "self-grooming actions such as hair pulling, nail biting, and skin picking "
    "that affect roughly two to five percent of the general population. That is "
    "millions of people worldwide.\n\n"
    "What makes these behaviors especially difficult to manage is that most people "
    "are completely unaware they are doing them. My system addresses this by using "
    "a standard webcam, computer vision through Google's MediaPipe library, and a "
    "machine learning classifier to detect these behaviors in real time and alert "
    "the user immediately.\n\n"
    "Over the next 25 minutes I will walk you through the full architecture, the "
    "database design, the machine learning pipeline, the detection engine, and "
    "then give a live demonstration of the working application."
)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2: AGENDA
# ═══════════════════════════════════════════════════════════════════════════
add_bullet_slide("Presentation Agenda", [
    "1.  Problem Statement: What are BFRBs and why detect them?",
    "2.  System Architecture: Tech stack and design decisions",
    "3.  Database Design: Schema, 3NF normalization, ERD",
    "4.  Backend Implementation: Flask API, session management",
    "5.  ML Pipeline: Random Forest classifier, feature engineering",
    "6.  Frontend UI: Real-time detection dashboard",
    "7.  Authentication & Onboarding: User management",
    "8.  Advanced Features: Calibration, bouts, pilot studies, PWA",
    "9.  Live Demonstration",
    "10. Results, Challenges & Future Work",
    "11. Q&A",
], "Here is our agenda for today. We will start with the problem statement — what "
   "BFRBs are and why detecting them matters clinically. Then I will walk through "
   "the full system architecture including the tech stack and design decisions I "
   "made along the way.\n\n"
   "Next, we will cover the database design — I have nine tables in Third Normal "
   "Form with proper foreign keys and constraints. After that, I will show you "
   "the backend implementation with over twenty Flask API endpoints covering "
   "authentication, sessions, detections, and analytics.\n\n"
   "Then comes the ML pipeline — how I trained a Random Forest classifier on "
   "Kaggle sensor data and deployed it for camera-based detection. We will look "
   "at the frontend dashboard, the authentication and onboarding system, and "
   "several advanced features including adaptive calibration, bout detection, "
   "pilot study mode, and Progressive Web App support.\n\n"
   "I will then do a live demonstration of the entire application — creating an "
   "account, going through onboarding, starting a detection session, and showing "
   "the analytics. We will wrap up with results, challenges, and future work, "
   "followed by Q and A. The whole presentation should take about 25 minutes.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3: PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════════
add_bullet_slide("Problem Statement: BFRBs", [
    "**What are BFRBs?**",
    "Body-Focused Repetitive Behaviors: compulsive self-grooming actions",
    "Hair pulling (trichotillomania), nail biting, skin picking, face touching",
    "Affects 2-5% of the population worldwide",
    "",
    "**The Challenge**",
    "Most people are UNAWARE when they're engaging in BFRBs",
    "Awareness is the first step to behavioral change",
    "Current solutions: rubber bands, journals, therapy alone",
    "",
    "**Our Solution**",
    "Real-time webcam detection using computer vision + ML",
    "Immediate visual alerts when BFRB is detected",
    "Track patterns over time for therapy support",
], "Let me start with the problem. Body-Focused Repetitive Behaviors, or BFRBs, "
   "are a class of compulsive self-grooming actions. The most common ones are "
   "trichotillomania — which is hair pulling — onychophagia, which is nail biting, "
   "and excoriation, which is skin picking. Other types include face touching, "
   "cheek biting, lip biting, nose picking, and eyelash pulling.\n\n"
   "These behaviors affect between two and five percent of the general population. "
   "To put that in perspective, that is roughly 15 to 40 million people in the "
   "United States alone. The key clinical challenge is that most people are "
   "completely unaware when they are engaging in these behaviors. You might be "
   "pulling your hair or biting your nails during a lecture and not realize it "
   "until you see the damage afterwards.\n\n"
   "Current approaches include rubber band snapping on the wrist, manual journaling, "
   "and therapy sessions — but none of these provide real-time awareness during the "
   "actual behavior.\n\n"
   "My solution is a web-based system that uses your laptop or desktop webcam to "
   "track your face and hands in real time using computer vision. When it detects "
   "that your hand is near your face in a pattern consistent with a BFRB — and "
   "sustains that contact for a sufficient duration — it alerts you immediately "
   "with a visual notification. Over time, it tracks your patterns so you and "
   "your therapist can identify trends and triggers.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4: SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════
add_two_col_slide("System Architecture",
    [
        "--- Frontend ---",
        "HTML5, CSS3, Vanilla JavaScript",
        "MediaPipe Holistic (face + hand tracking)",
        "Chart.js (statistical dashboard)",
        "Socket.IO client (real-time)",
        "Progressive Web App (PWA)",
        "",
        "--- Computer Vision ---",
        "468 face landmarks per frame",
        "21 hand landmarks per frame",
        "One Euro Filter for smoothing",
        "13 anatomical face zones",
    ],
    [
        "--- Backend ---",
        "Python 3.11 + Flask 2.3",
        "SQLite (embedded database)",
        "Flask-SocketIO (WebSocket)",
        "scikit-learn (ML inference)",
        "",
        "--- ML Pipeline ---",
        "Random Forest classifier",
        "16 semantic features",
        "Trained on Kaggle CMI dataset",
        "574,945 readings, 8,151 sequences",
        "Late multimodal fusion",
    ],
    "Now let me walk you through the system architecture. The application has two main "
    "layers: the frontend running in the browser, and the backend running on a Python "
    "Flask server.\n\n"
    "On the frontend side, I use plain HTML5, CSS3, and vanilla JavaScript — no React "
    "or Angular, keeping dependencies minimal. The computer vision runs entirely in the "
    "browser using Google's MediaPipe Holistic library, which gives me 468 face landmarks "
    "and 21 hand landmarks per hand, per frame, in real time. This is important for "
    "privacy — the video feed never leaves the user's device. All landmark extraction "
    "happens client-side.\n\n"
    "I use Chart.js for the statistical dashboard with four different chart types: "
    "daily detection trends, per-type distribution, session timeline, and accuracy "
    "tracking. Socket.IO provides a persistent WebSocket connection for low-latency "
    "ML predictions. And the entire app is a Progressive Web App, meaning it has a "
    "manifest, a service worker, and can be installed on the desktop.\n\n"
    "On the backend, I use Python 3.11 with Flask 2.3 as the web framework. SQLite "
    "is the embedded database — zero configuration, no Docker needed. Flask-SocketIO "
    "handles the WebSocket layer for real-time ML inference, and scikit-learn provides "
    "the Random Forest classifier.\n\n"
    "The ML pipeline was trained on the Kaggle CMI 'Detect Behavior with Sensor Data' "
    "dataset — 574,945 readings from 80 participants wearing wrist sensors across "
    "8,151 labeled sequences. I extract 16 semantic features that bridge the gap "
    "between wrist sensor measurements and camera-derived landmarks, and use late "
    "multimodal fusion to combine the ML classifier output with the rule-based "
    "detection engine.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5: DATABASE DESIGN
# ═══════════════════════════════════════════════════════════════════════════
add_bullet_slide("Database Design (3NF Normalized)", [
    "**9 Tables in Third Normal Form (3NF)**",
    "users - accounts with hashed passwords (SHA-256 + salt)",
    "onboarding_responses - personalization questionnaire data",
    "bfrb_types - 8 seeded BFRB categories (extensible)",
    "sessions - detection sessions (start, end, stats)",
    "detections - individual events (confidence, ML score, fused score)",
    "calibrations - per-session adaptive thresholds",
    "user_feedback - yes/no feedback on each detection",
    "pilot_trials - structured ground-truth collection",
    "telemetry - per-frame IEEE-grade logging",
    "",
    "**Key Design Decisions**",
    "SQLite for zero-config deployment (no Docker needed)",
    "Foreign keys + WAL journaling for data integrity",
    "Auto-migration: new columns added transparently",
], "The database is designed in Third Normal Form — 3NF — which means every "
   "non-key attribute depends on the key, the whole key, and nothing but the key. "
   "There are no transitive dependencies and no repeating groups.\n\n"
   "I have nine tables total. The users table stores accounts with SHA-256 salted "
   "password hashes — I never store plaintext passwords. The onboarding_responses "
   "table stores the personalization questionnaire answers, linked by a foreign "
   "key to the users table.\n\n"
   "The bfrb_types table is seeded with eight standard BFRB categories but is "
   "extensible — users can add custom types. Sessions tracks each detection "
   "session with start and end timestamps. Detections stores every individual "
   "BFRB event with its confidence score, ML score, fused score, and the zone "
   "where the touch was detected.\n\n"
   "Calibrations stores the per-session adaptive thresholds computed during "
   "the calibration wizard. User_feedback captures yes or no responses on each "
   "detection — this is used for accuracy tracking. Pilot_trials supports "
   "structured ground-truth collection for research. And telemetry stores "
   "per-frame IEEE-grade logs for full reproducibility.\n\n"
   "I chose SQLite for the embedded database because it requires zero "
   "configuration — no separate database server, no Docker. The database file "
   "is created automatically on first run. I also provide a MySQL reference "
   "schema in the database directory for anyone who wants to deploy to "
   "production with a full RDBMS. WAL journaling is enabled for concurrent "
   "read and write safety.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6: BACKEND API
# ═══════════════════════════════════════════════════════════════════════════
add_two_col_slide("Backend: Flask API (20+ Endpoints)",
    [
        "--- Authentication ---",
        "POST /api/auth/register",
        "POST /api/auth/login",
        "POST /api/auth/logout",
        "GET  /api/auth/me",
        "",
        "--- Sessions ---",
        "POST /api/sessions (start)",
        "PUT  /api/sessions/<id>/end",
        "GET  /api/sessions (list)",
        "",
        "--- Detection & Feedback ---",
        "POST /api/detections",
        "POST /api/feedback",
    ],
    [
        "--- ML Inference ---",
        "POST /api/predict (HTTP)",
        "WS predict_frame (WebSocket)",
        "POST /api/calibrate",
        "",
        "--- Analytics ---",
        "GET /api/stats/overview",
        "GET /api/stats/daily",
        "GET /api/stats/trends",
        "GET /api/bouts/<sid>",
        "",
        "--- Research ---",
        "POST /api/pilot/start",
        "GET /api/export/session/<sid>",
    ],
    "The backend exposes over twenty RESTful API endpoints organized into six "
    "categories. All endpoints follow standard HTTP method conventions — POST for "
    "creation, GET for retrieval, PUT for updates — and return consistent JSON "
    "responses with proper HTTP status codes.\n\n"
    "Authentication has four endpoints: register creates a new account with "
    "password hashing, login validates credentials and creates a Flask session, "
    "logout destroys the session, and the me endpoint returns the currently "
    "logged-in user's profile.\n\n"
    "Session management lets users start a detection session with POST, end it "
    "with PUT, and list all their sessions with GET. Each session tracks total "
    "detections, duration, and the BFRB types that were monitored.\n\n"
    "The detection and feedback endpoints log individual BFRB events and capture "
    "user confirmation — did you actually just bite your nails, yes or no? This "
    "feedback loop is critical for measuring real-world accuracy.\n\n"
    "For ML inference, there are two pathways: a standard HTTP POST endpoint "
    "at /api/predict, and a WebSocket event called predict_frame. The WebSocket "
    "path is preferred because it avoids HTTP overhead and delivers predictions "
    "with about 5 to 15 milliseconds of round-trip latency. The calibrate "
    "endpoint saves adaptive threshold parameters.\n\n"
    "Analytics endpoints provide lifetime statistics, daily detection counts, "
    "trend analysis, and bout aggregation. The research endpoints support "
    "pilot studies with structured ground-truth collection and CSV export "
    "for offline analysis in tools like Excel or R.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7: ML PIPELINE
# ═══════════════════════════════════════════════════════════════════════════
add_bullet_slide("ML Pipeline: Feature Engineering & Classification", [
    "**Training Data: Kaggle CMI 'Detect Behavior with Sensor Data'**",
    "80 participants wearing wrist sensors | 574,945 readings",
    "18 gesture types (8 BFRB-like + 10 normal controls)",
    "",
    "**16 Semantic Features (bridging sensors -> camera)**",
    "Accelerometer -> wrist velocity (mean, std, max)",
    "Gyroscope -> palm rotation dynamics",
    "Thermopile -> hand-face overlap fraction",
    "Time-of-Flight -> 5 anatomical distances (nose, mouth, cheek, forehead, chin)",
    "",
    "**Model: Random Forest (200 estimators)**",
    "5-fold cross-validation | Macro F1 evaluation",
    "MinMaxScaler normalization [0, 1]",
    "",
    "**Late Multimodal Fusion**",
    "fused_score = 0.6 x ML_score + 0.4 x rule_score",
], "This is one of the most technically interesting parts of the project — the "
   "machine learning pipeline.\n\n"
   "The training data comes from the Kaggle CMI competition called 'Detect "
   "Behavior with Sensor Data.' This dataset has 80 participants wearing wrist "
   "sensors — accelerometers, gyroscopes, thermopiles, and time-of-flight "
   "sensors — performing 18 different gesture types, including 8 that are "
   "BFRB-like behaviors and 10 normal control gestures.\n\n"
   "Now here is the key challenge: the training data is from wrist sensors, "
   "but my system uses a camera. So I designed 16 semantic features that "
   "bridge these two modalities. For example, the accelerometer readings "
   "from the wrist sensor correspond to the wrist velocity I can compute "
   "from hand landmarks — mean velocity, standard deviation, and max velocity. "
   "The gyroscope readings map to palm rotation dynamics that I derive from "
   "the orientation of the hand landmarks. The thermopile — which measures "
   "infrared heat — corresponds to the hand-face overlap fraction I compute "
   "geometrically. And the time-of-flight sensor measurements map to five "
   "anatomical distances I calculate between the fingertip and face zones: "
   "nose, mouth, cheek, forehead, and chin.\n\n"
   "The model itself is a Random Forest with 200 decision trees, evaluated "
   "using 5-fold cross-validation with Macro F1 as the metric. I use "
   "MinMaxScaler to normalize all features to the 0-to-1 range.\n\n"
   "Finally, I use late multimodal fusion to combine the ML classifier output "
   "with the rule-based detection score. The formula is: fused score equals "
   "0.6 times the ML score plus 0.4 times the rule-based score. This way, "
   "both systems have to agree before a detection is triggered, which "
   "significantly reduces false positives.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8: DETECTION PIPELINE
# ═══════════════════════════════════════════════════════════════════════════
add_bullet_slide("Detection Pipeline (9-Stage)", [
    "1. One Euro Filter - adaptive low-pass smoothing on all landmarks",
    "2. Multi-Zone Face Detection - 13 anatomical regions with centroids",
    "3. Hand Pose Classifier - pinch_grip / fist / pointing / open_hand / curled",
    "4. 4-State Rule Engine - IDLE -> PROXIMITY -> DWELL -> BFRB",
    "5. Multi-Signal BFRB Score - 11 weighted factors (proximity, dwell, curl...)",
    "6. ML Classifier - 16-feature Random Forest, 1.5s sliding window",
    "7. Late Fusion - combines ML + rule scores (0.6/0.4 split)",
    "8. Adaptive Threshold - per-session calibration (baseline vs touch distance)",
    "9. IEEE Telemetry - per-frame structured logging for reproducibility",
    "",
    "**False Positive Filters**",
    "Eating detection (mouth wide open + hand near mouth + open palm)",
    "Resting detection (static hand + open palm + no recent BFRB history)",
], "The detection pipeline has nine stages that run every single frame — about "
   "25 to 30 times per second.\n\n"
   "Stage one is the One Euro Filter, an adaptive low-pass filter that smooths "
   "the raw landmark positions. Without this, small jittery movements in the "
   "webcam feed would cause noisy detections. The filter adapts its cutoff "
   "frequency based on the speed of movement — slow movements get heavy "
   "smoothing, fast movements pass through quickly.\n\n"
   "Stage two is multi-zone face detection. I divide the face into 13 "
   "anatomical regions — forehead, left temple, right temple, left eye, "
   "right eye, nose bridge, nose tip, left cheek, right cheek, upper lip, "
   "lower lip, chin, and jawline. Each zone is defined by a cluster of "
   "MediaPipe face landmark indices with a computed centroid.\n\n"
   "Stage three classifies the hand pose: is it a pinch grip, a fist, "
   "pointing, open hand, or curled fingers? This matters because a pinch "
   "grip near the hairline strongly suggests hair pulling, while an open "
   "palm near the mouth is more likely eating.\n\n"
   "Stage four is the 4-state rule engine — a state machine that transitions "
   "through IDLE, PROXIMITY, DWELL, and BFRB. The hand must first enter the "
   "proximity zone, then dwell there for 1.2 seconds, and then sustain "
   "contact for another 2.5 seconds before a BFRB is triggered. This "
   "prevents momentary face touches from causing false alarms.\n\n"
   "Stage five computes a multi-signal BFRB score from 11 weighted factors "
   "including proximity, dwell time, finger curl, velocity, and more. "
   "Stage six sends 16 features to the server-side Random Forest for ML "
   "classification. Stage seven fuses the ML and rule scores. Stage eight "
   "applies the adaptive threshold from calibration. And stage nine logs "
   "everything to the telemetry table for reproducibility.\n\n"
   "I also have false positive filters for eating — detected by mouth wide "
   "open plus hand near mouth plus open palm — and for resting, which is "
   "a static hand with an open palm and no recent BFRB history. These "
   "are NOT BFRBs even though the hand is near the face.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9: FRONTEND UI
# ═══════════════════════════════════════════════════════════════════════════
add_bullet_slide("Frontend: Professional Web Dashboard", [
    "**Clean White + Dark Blue Professional Theme**",
    "Responsive design (desktop + tablet)",
    "Real-time video feed with landmark overlays",
    "",
    "**Three Main Views**",
    "Dashboard: live detection feed + controls + BFRB selector + signals",
    "Session HUD: charts (daily detections, type distribution, timeline, accuracy)",
    "History: session table with stats per session",
    "",
    "**Interactive Elements**",
    "BFRB type multi-select (8 types, multi/single mode)",
    "Detection notifications with yes/no feedback buttons",
    "Calibration wizard (2-step: far baseline + touch measurement)",
    "Pilot study modal (structured ground-truth collection)",
    "CSV export for offline analysis",
    "",
    "**Authentication: sign up, sign in, onboarding questionnaire**",
], "The frontend is a single-page application with a clean, professional design. "
   "I deliberately chose a white background with dark navy blue accents rather "
   "than a flashy dark AI-looking theme. This is a health and wellness tool — "
   "it needs to feel trustworthy and calm, not like a sci-fi dashboard.\n\n"
   "The interface has three main views, all navigated by tabs at the top. The "
   "Dashboard view shows the live webcam feed with landmark overlays drawn in "
   "real time, detection controls, the BFRB type selector where you can choose "
   "which behaviors to monitor, and real-time signal indicators showing the "
   "current state of the detection engine.\n\n"
   "The Session HUD view has four charts: a line chart of daily detection counts, "
   "a doughnut chart showing the distribution of BFRB types, a timeline scatter "
   "plot of individual detections within the current session, and an accuracy "
   "chart based on user feedback.\n\n"
   "The History view shows a table of all past sessions with statistics like "
   "total detections, session duration, and which types were detected.\n\n"
   "Interactive elements include the BFRB type selector with multi-select "
   "chips for eight types, detection notification cards with yes and no "
   "feedback buttons, the calibration wizard which walks you through a "
   "two-step process, the pilot study modal for structured research data "
   "collection, and CSV export for downloading session data.\n\n"
   "The authentication system includes a sign-up and sign-in page, and "
   "first-time users go through a six-question onboarding questionnaire "
   "that personalizes which BFRB types are monitored by default.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10: AUTH & ONBOARDING
# ═══════════════════════════════════════════════════════════════════════════
add_bullet_slide("Authentication & User Onboarding", [
    "**Secure Authentication**",
    "Password hashing: SHA-256 with random 16-byte salt",
    "Flask session-based auth with login_required decorator",
    "Username or email login support",
    "",
    "**Onboarding Questionnaire (first login)**",
    "Which BFRBs do you experience? (multi-select)",
    "How often do these behaviors occur? (5 frequency levels)",
    "What triggers your BFRBs? (stress, boredom, anxiety, focus...)",
    "What is your primary goal? (awareness, reduce, track, research, clinical)",
    "Alert preferences (visual only, visual+sound, gentle, silent)",
    "",
    "**Personalization**",
    "Pre-selects the user's reported BFRB types in the detector",
    "Responses stored in onboarding_responses table for future use",
], "Security was a priority even though this is an academic project. Passwords "
   "are never stored in plaintext. When a user registers, I generate a random "
   "16-byte salt, concatenate it with the password, and hash the result using "
   "SHA-256. Both the salt and the hash are stored in the database. During "
   "login, I regenerate the hash from the submitted password and the stored "
   "salt, then compare. This is the same approach used in production systems.\n\n"
   "Authentication is session-based using Flask's built-in session management. "
   "I wrote a login_required decorator that protects all API endpoints — if "
   "you are not logged in, you get redirected to the login page. Users can "
   "log in with either their username or email address.\n\n"
   "When a new user logs in for the first time, they see an onboarding "
   "questionnaire overlay with six questions. The first asks which BFRBs "
   "they experience — this is a multi-select with clickable chips for all "
   "eight types. The second asks how frequently these behaviors occur, from "
   "rarely to constantly. The third asks about triggers — stress, boredom, "
   "anxiety, focus, fatigue, and others.\n\n"
   "The fourth question asks about their primary goal: is it awareness, "
   "reduction, tracking, research, or clinical use? The fifth asks about "
   "alert preferences — visual only, visual plus sound, gentle, or silent. "
   "And the sixth is an optional free-text field for additional notes.\n\n"
   "All responses are saved to the onboarding_responses table and used to "
   "pre-select the user's reported BFRB types in the detector, so the system "
   "is personalized from the very first session. Users can skip the "
   "onboarding if they prefer and set things up manually later.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11: ADVANCED FEATURES
# ═══════════════════════════════════════════════════════════════════════════
add_bullet_slide("Advanced Features", [
    "**Adaptive Threshold Calibration (Tier 2)**",
    "2-step wizard: measure far-baseline and touch distance",
    "threshold = touch + 0.3 * (baseline - touch)",
    "Accounts for different face sizes and camera distances",
    "",
    "**Bout Detection (Tier 3)**",
    "Groups consecutive detections into clinical 'bouts'",
    "5-second gap threshold | per-type aggregation",
    "",
    "**Pilot Study Mode (Tier 3)**",
    "Structured ground-truth collection: participant, gesture, duration",
    "Computes precision, recall, F1 per trial",
    "",
    "**WebSocket Real-Time (Tier 3)**",
    "Flask-SocketIO: ML predictions over persistent socket",
    "250ms throttle on WS vs 400ms on HTTP fallback",
    "",
    "**PWA (Tier 4) - Offline-capable app shell**",
], "Let me walk through the advanced features that take this project beyond "
   "a basic detection system.\n\n"
   "First, adaptive threshold calibration. Different users sit at different "
   "distances from their camera, and different people have different face "
   "sizes. A fixed proximity threshold would not work for everyone. So I "
   "built a two-step calibration wizard. In step one, you sit normally and "
   "hold still — the system measures your baseline hand-to-face distance. "
   "In step two, you touch your face — the system measures the touch "
   "distance. The adaptive threshold is then computed as: threshold equals "
   "the touch distance plus alpha times the difference between baseline "
   "and touch, where alpha is 0.3. This normalizes the detection for each "
   "individual user and camera setup.\n\n"
   "Second, bout detection. In clinical BFRB research, individual detections "
   "are less meaningful than 'bouts' — sustained episodes of repetitive "
   "behavior. My system groups consecutive detections that occur within "
   "a 5-second gap into a single bout, and aggregates them per BFRB type. "
   "This gives clinicians a more meaningful metric.\n\n"
   "Third, pilot study mode. This is a structured data collection interface "
   "for ground-truth validation. You specify a participant ID, a gesture "
   "type, and a duration, and the system records exactly what happened "
   "during that window. It then computes precision, recall, and F1 score "
   "for that trial.\n\n"
   "Fourth, WebSocket real-time inference. Instead of making an HTTP POST "
   "request for every frame, the client sends landmark data over a "
   "persistent Socket.IO connection. The server runs the Random Forest "
   "classifier and sends back the prediction — all within about 5 to 15 "
   "milliseconds. The client throttles to one prediction every 250 "
   "milliseconds to avoid overloading the server.\n\n"
   "And finally, the app is a Progressive Web App with a manifest file "
   "and a service worker that caches the app shell for offline use.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12: LIVE DEMO PLAN
# ═══════════════════════════════════════════════════════════════════════════
add_bullet_slide("Live Demonstration", [
    "**Demo Flow (5-7 minutes)**",
    "",
    "1. Login Page - create a new account",
    "2. Onboarding - answer setup questionnaire",
    "3. Dashboard - show the clean white + dark blue interface",
    "4. Start Session - camera activates, MediaPipe loads",
    "5. Show Detection - bring hand near face, trigger BFRB detection",
    "6. Feedback - respond to detection notification (yes/no)",
    "7. Calibration - run 2-step adaptive threshold calibration",
    "8. Session HUD - switch to charts tab, show daily/type/timeline charts",
    "9. History - view past sessions in the table",
    "10. Export CSV - download session data",
    "11. Stop Session - end and show summary",
    "12. Sign Out",
], "Now I am going to do a live demonstration. Let me switch to the browser.\n\n"
   "Step one — I will open the login page. You can see the clean white and "
   "dark blue design. I will create a brand new account right now with a "
   "username, email, full name, and password.\n\n"
   "Step two — after account creation and login, the onboarding questionnaire "
   "appears. I will select a few BFRB types — let us say hair pulling and "
   "nail biting — set the frequency, choose some triggers, pick my goal, "
   "and set alert preferences. Then I will submit.\n\n"
   "Step three — now we are on the main dashboard. Notice the professional "
   "white and navy theme. The webcam feed is on the left, detection controls "
   "on the right. The BFRB types I selected during onboarding are already "
   "pre-checked.\n\n"
   "Step four — I will click Start Session. The camera activates and "
   "MediaPipe starts loading. You should see green face landmarks and "
   "blue hand landmarks appearing on the video feed.\n\n"
   "Step five — now I will bring my hand near my face to demonstrate "
   "detection. Watch the state indicator — it should go from IDLE to "
   "PROXIMITY as my hand approaches, then to DWELL after about 1.2 "
   "seconds of sustained contact, and finally to BFRB after about 2.5 "
   "seconds total. A notification card will appear.\n\n"
   "Step six — I will respond to the notification with yes or no feedback. "
   "Step seven — let me run the calibration wizard. Step eight — switching "
   "to the Session HUD tab to show the charts. Step nine — the History "
   "tab shows past sessions. Step ten — I will export a CSV. Step eleven — "
   "stop the session. Step twelve — sign out.\n\n"
   "Let me go ahead and do this now.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13: GITHUB REPOSITORY
# ═══════════════════════════════════════════════════════════════════════════
add_bullet_slide("GitHub Repository & Deployment", [
    "**Repository: github.com/Nishanth3424/bfrb-detection-eneb453**",
    "",
    "**Project Structure**",
    "bfrb_web_app/app.py         - Flask backend (1,300+ lines)",
    "bfrb_web_app/templates/     - index.html (3,300+ lines) + login.html",
    "bfrb_web_app/model/         - Trained ML model + feature metadata",
    "bfrb_web_app/static/        - PWA manifest + service worker + icons",
    "database/                   - MySQL schema, seed data, demo queries",
    "erd/                        - Interactive ERD visualization",
    "wireframe/                  - UI wireframe mockup",
    "",
    "**How to Run**",
    "pip install -r requirements.txt",
    "python app.py",
    "Open http://localhost:5000",
    "",
    "**Total: ~4,800 lines of code (Python + HTML/CSS/JS)**",
], "The entire project is hosted on GitHub as a public repository. You can "
   "clone it and run it locally on any machine with Python 3.11 installed.\n\n"
   "Let me walk through the project structure. The main application code "
   "lives in the bfrb_web_app directory. The app.py file is the Flask "
   "backend — over 1,300 lines of Python covering all the API endpoints, "
   "database initialization, ML inference, and WebSocket handling.\n\n"
   "The templates directory contains two HTML files: index.html is the "
   "main web app at over 3,300 lines — it contains all the frontend "
   "JavaScript, CSS, and HTML in a single file for simplicity. Login.html "
   "is the authentication page.\n\n"
   "The model directory contains the trained Random Forest pickle file, "
   "the MinMaxScaler pickle, and a feature metadata JSON file that "
   "documents exactly which 16 features the model expects.\n\n"
   "The static directory has the PWA manifest, the service worker, and "
   "app icons. The database directory contains the MySQL reference schema, "
   "seed data with sample records, and demo SQL queries. The erd directory "
   "has an interactive entity-relationship diagram as an HTML file. And "
   "the wireframe directory has the original UI mockup.\n\n"
   "To run the project, you just need three commands: clone the repo, "
   "pip install the requirements, and python app.py. It starts on "
   "localhost port 5000. No Docker, no database server setup, no "
   "environment variables to configure. SQLite creates the database "
   "file automatically on first run.\n\n"
   "In total, the project is approximately 4,800 lines of code across "
   "Python, HTML, CSS, and JavaScript.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14: RESULTS & METRICS
# ═══════════════════════════════════════════════════════════════════════════
add_bullet_slide("Results & Performance", [
    "**ML Model Performance**",
    "Random Forest: 77% accuracy (5-fold CV)",
    "Macro F1: handles class imbalance (64% BFRB / 36% non-BFRB)",
    "Training on 8,151 labeled sequences from 80 participants",
    "",
    "**Real-Time Performance**",
    "25-30 FPS with MediaPipe Holistic in browser",
    "ML inference: <10ms per prediction (server-side)",
    "WebSocket latency: ~5-15ms round-trip",
    "",
    "**Detection Capability**",
    "8 BFRB types supported (+ custom types)",
    "Multi-zone face detection (13 anatomical regions)",
    "False positive filtering (eating, resting)",
    "Adaptive calibration per user",
    "",
    "**Code Coverage: 20+ API endpoints, 9 database tables**",
], "Let me share the key results and performance metrics.\n\n"
   "The Random Forest ML model achieves 77 percent accuracy under 5-fold "
   "cross-validation on the Kaggle CMI dataset. I use Macro F1 as the "
   "primary evaluation metric because the dataset has class imbalance — "
   "about 64 percent of the labeled sequences are BFRB behaviors and "
   "36 percent are non-BFRB controls. Macro F1 gives equal weight to "
   "each class regardless of size, so it is a fairer measure than "
   "raw accuracy.\n\n"
   "Now, 77 percent might seem moderate, but consider the context: I am "
   "training on wrist sensor data and deploying on camera-derived features. "
   "This cross-modal domain transfer is inherently lossy. The fact that "
   "the semantic feature mapping preserves enough signal to achieve 77 "
   "percent is actually quite strong.\n\n"
   "For real-time performance, MediaPipe Holistic runs at 25 to 30 frames "
   "per second in the browser, which is smooth enough for continuous "
   "monitoring. The ML inference on the server side takes less than 10 "
   "milliseconds per prediction — Random Forests are very fast at "
   "inference time. And the WebSocket round-trip latency is about 5 to "
   "15 milliseconds, so the user gets near-instant predictions.\n\n"
   "The system supports all eight standard BFRB types plus custom types, "
   "detects across 13 anatomical face zones, filters out false positives "
   "from eating and resting, and adapts to each user through calibration.\n\n"
   "From a code coverage perspective, there are over 20 API endpoints "
   "and 9 database tables, all properly normalized and tested.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15: CHALLENGES
# ═══════════════════════════════════════════════════════════════════════════
add_bullet_slide("Challenges & Solutions", [
    "**Challenge 1: Sensor-to-Camera Domain Gap**",
    "Kaggle data = wrist sensors, our system = camera",
    "Solution: semantic feature mapping (16 features bridge both modalities)",
    "",
    "**Challenge 2: False Positives**",
    "Eating, resting, scratching look similar to BFRBs",
    "Solution: multi-signal scoring + eating/resting filters + dwell timers",
    "",
    "**Challenge 3: Varying User Distances**",
    "Proximity thresholds break at different camera distances",
    "Solution: adaptive calibration normalizes by face width",
    "",
    "**Challenge 4: Real-Time Performance**",
    "MediaPipe + ML inference must run at >20 FPS",
    "Solution: One Euro filter + sliding window + WebSocket transport",
    "",
    "**Challenge 5: User Experience**",
    "Technical ML system must be accessible to non-technical users",
    "Solution: onboarding questionnaire + clean professional UI",
], "I want to be transparent about the challenges I faced and how I solved them.\n\n"
   "The biggest challenge was the sensor-to-camera domain gap. The Kaggle "
   "training data comes from wrist-worn sensors — accelerometers, gyroscopes, "
   "thermopiles, and time-of-flight sensors. But my system uses a webcam. "
   "These are fundamentally different measurement modalities. My solution "
   "was to design 16 semantic features that capture the same physical "
   "quantities — velocity, rotation, proximity, distance — from camera "
   "landmarks instead of sensor readings. This semantic mapping is what "
   "makes the cross-modal transfer possible.\n\n"
   "The second challenge was false positives. Eating lunch, resting your "
   "chin on your hand, and scratching an itch all involve your hand near "
   "your face — but they are not BFRBs. I addressed this with multiple "
   "layers: a multi-signal scoring system that requires several indicators "
   "to agree, specific eating and resting filters that check for open "
   "palms and mouth position, and dwell timers that require sustained "
   "contact for at least 2.5 seconds before triggering.\n\n"
   "The third challenge was handling different user distances from the "
   "camera. If you sit two feet away versus four feet away, the raw "
   "landmark distances are completely different. The adaptive calibration "
   "wizard normalizes everything relative to your specific face width "
   "and camera distance.\n\n"
   "The fourth challenge was real-time performance. Running MediaPipe "
   "landmark extraction plus ML inference plus UI rendering at over 20 "
   "FPS requires careful optimization. The One Euro filter reduces noise "
   "without adding latency, the sliding window batches features efficiently, "
   "and the WebSocket transport eliminates HTTP overhead.\n\n"
   "The fifth challenge was making a technical ML system accessible to "
   "non-technical users. The onboarding questionnaire, clean professional "
   "UI, and step-by-step calibration wizard all address this.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16: FUTURE WORK
# ═══════════════════════════════════════════════════════════════════════════
add_bullet_slide("Future Work", [
    "**Short-Term Improvements**",
    "Fine-tune ML model on camera-derived features (labeled webcam data)",
    "Add multi-user support with role-based access (patient vs therapist)",
    "Mobile-optimized layout for smartphone webcam use",
    "",
    "**Research Extensions**",
    "Federated learning: train across users without sharing raw data",
    "CORAL domain adaptation: align sensor and camera feature distributions",
    "Longitudinal studies: track behavior patterns over weeks/months",
    "",
    "**Clinical Integration**",
    "Therapist dashboard: view patient BFRB trends remotely",
    "Integration with CBT (Cognitive Behavioral Therapy) protocols",
    "Smart notifications: context-aware alerts based on time/location",
    "",
    "**Deployment**",
    "Cloud hosting (AWS/GCP) for multi-user access",
    "Browser extension version for always-on monitoring",
], "Looking ahead, there are several directions this project could go.\n\n"
   "In the short term, the most impactful improvement would be fine-tuning "
   "the ML model on camera-derived features — collecting labeled webcam "
   "data from real BFRB sessions and retraining. This would close the "
   "domain gap between sensor and camera and likely push accuracy above "
   "85 percent. I would also add multi-user support with role-based "
   "access, so a therapist could view their patients' data, and optimize "
   "the layout for mobile devices.\n\n"
   "For research extensions, federated learning would allow the model to "
   "improve across many users without any raw video data leaving their "
   "devices — critical for healthcare privacy. CORAL domain adaptation "
   "is a technique that aligns the statistical distributions of sensor "
   "features and camera features, which could further improve the "
   "cross-modal transfer. And longitudinal studies tracking behavior "
   "patterns over weeks or months would provide valuable clinical insights.\n\n"
   "For clinical integration, a therapist-facing dashboard would let "
   "clinicians monitor patient BFRB trends remotely between sessions. "
   "Integration with Cognitive Behavioral Therapy protocols could provide "
   "structured habit reversal training guided by the detection data. "
   "And context-aware smart notifications could adjust alert behavior "
   "based on time of day or user activity.\n\n"
   "For deployment, moving to a cloud-hosted solution on AWS or GCP would "
   "enable multi-user access without local installation, and a browser "
   "extension version could provide always-on monitoring during daily "
   "computer use.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17: REQUIREMENTS CHECKLIST
# ═══════════════════════════════════════════════════════════════════════════
add_two_col_slide("Project Requirements Checklist",
    [
        "--- Database ---",
        "[DONE] 3NF normalized schema",
        "[DONE] Proper foreign keys & indexes",
        "[DONE] 9 tables (users, sessions, etc.)",
        "[DONE] MySQL reference schema + Docker",
        "[DONE] ERD diagram",
        "",
        "--- Backend ---",
        "[DONE] Flask routes (20+ endpoints)",
        "[DONE] Session management",
        "[DONE] Error handling (JSON responses)",
        "[DONE] ML inference endpoint",
        "[DONE] WebSocket real-time",
    ],
    [
        "--- Frontend ---",
        "[DONE] Responsive HTML/CSS/JS",
        "[DONE] Real-time webcam detection",
        "[DONE] Charts & analytics dashboard",
        "[DONE] Authentication (login/signup)",
        "[DONE] Onboarding questionnaire",
        "",
        "--- Extras ---",
        "[DONE] BFRB type selector (custom)",
        "[DONE] Adaptive calibration",
        "[DONE] CSV export",
        "[DONE] PWA (offline capable)",
        "[DONE] GitHub repository",
    ],
    "Before we wrap up, let me walk through the project requirements checklist "
    "to confirm everything is covered.\n\n"
    "For the database: we have a fully normalized 3NF schema with 9 tables, "
    "proper foreign keys and indexes on all lookup columns, a MySQL reference "
    "schema for production deployment, and an interactive ERD diagram. All done.\n\n"
    "For the backend: over 20 Flask API endpoints following RESTful conventions, "
    "session-based authentication with proper security, consistent JSON error "
    "handling across all endpoints, ML inference via both HTTP and WebSocket, "
    "and real-time communication through Flask-SocketIO. All done.\n\n"
    "For the frontend: a responsive HTML, CSS, and JavaScript interface that "
    "works on desktop and tablet, real-time webcam detection with MediaPipe "
    "landmark visualization, four different charts on the analytics dashboard, "
    "a complete authentication flow with sign-up and sign-in pages, and an "
    "onboarding questionnaire for new users. All done.\n\n"
    "And for the extras that go beyond the basic requirements: a custom BFRB "
    "type selector with multi-select support, adaptive threshold calibration "
    "that personalizes the detection for each user, CSV export for offline "
    "analysis, Progressive Web App support with offline caching, and a "
    "fully public GitHub repository with documentation. All done.\n\n"
    "Every single requirement in the rubric has been met or exceeded.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18: THANK YOU & Q&A
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY_DARK)
add_rect(slide, Inches(0), Inches(3.2), SLIDE_W, Inches(0.06), ACCENT_RED)
add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
             "Thank You!", size=52, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.6),
             "Questions & Discussion", size=28, color=GOLD,
             align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             "Nishanth S.  |  ENEB453  |  University of Maryland",
             size=18, color=TEXT_DIM, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
             "GitHub: github.com/Nishanth3424/bfrb-detection-eneb453",
             size=16, color=BLUE_LT, align=PP_ALIGN.CENTER)
slide.notes_slide.notes_text_frame.text = (
    "Thank you very much for your time and attention. I am happy to take any "
    "questions you might have about the project.\n\n"
    "Here are some questions I anticipate and my prepared answers:\n\n"
    "Q: Why did you choose MediaPipe over other computer vision libraries like "
    "OpenCV or TensorFlow.js?\n"
    "A: MediaPipe Holistic is free, runs natively in the browser with no server-side "
    "video processing, provides both face and hand landmarks in a single model, and "
    "achieves real-time performance at 25-30 FPS. OpenCV would require server-side "
    "processing which raises privacy concerns. TensorFlow.js could work but MediaPipe "
    "is purpose-built for landmark detection and is more efficient.\n\n"
    "Q: Why Random Forest instead of a deep learning model like an LSTM or CNN?\n"
    "A: The 16 features I extract are tabular, not sequential or spatial. Random "
    "Forest is ideal for tabular data — it is fast at inference time (under 10ms), "
    "interpretable (you can inspect feature importances), and does not require a "
    "GPU. For 16 numeric features, it outperforms neural networks in practice.\n\n"
    "Q: How accurate is the real-world detection?\n"
    "A: The ML model achieves 77 percent on the Kaggle validation set. Real-world "
    "accuracy depends on calibration, lighting, and camera quality. The late fusion "
    "approach and dwell timers significantly reduce false positives compared to "
    "using either the ML model or rule engine alone.\n\n"
    "Q: What about privacy? Is the video being sent to a server?\n"
    "A: No. All video processing and landmark extraction happens entirely in the "
    "browser using MediaPipe. The video feed never leaves the user's device. Only "
    "the 16 numeric features — numbers, not images — are sent to the server for "
    "ML classification.\n\n"
    "Q: Could this work on a phone?\n"
    "A: The PWA can be installed on a phone, and MediaPipe does support mobile "
    "browsers. The current UI layout is optimized for desktop, but making it "
    "mobile-responsive is a straightforward future improvement.\n\n"
    "Thank you again. The GitHub repository is linked on screen if you would "
    "like to try it yourself."
)

# ── Save ──
prs.save(OUT)
print(f"Presentation saved to: {OUT}")
print(f"Total slides: {len(prs.slides)}")
